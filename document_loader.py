import os
from typing import List, Dict, Optional
import docx2txt
from PyPDF2 import PdfReader
from pptx import Presentation
from config import DATA_DIR

class DocumentLoader:
    def __init__(
        self,
        data_dir: str = DATA_DIR,
    ):
        self.data_dir = data_dir
        self.supported_formats = [".pdf", ".pptx", ".docx", ".txt"]

    def load_pdf(self, file_path: str) -> List[Dict]:
        """加载PDF文件，按页返回内容"""
        try:
            reader = PdfReader(file_path)
            pages = []
            
            for page_num, page in enumerate(reader.pages):
                text = page.extract_text()
                if text.strip():  # 只添加有内容的页面
                    pages.append({
                        "text": f"--- 第 {page_num + 1} 页 ---\n{text}\n",
                        "page_number": page_num + 1
                    })
            
            return pages
        except Exception as e:
            print(f"加载PDF文件出错: {file_path}, 错误: {str(e)}")
            return []

    def load_pptx(self, file_path: str) -> List[Dict]:
        """加载PPT文件，按幻灯片返回内容"""
        try:
            prs = Presentation(file_path)
            slides = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                
                # 提取所有形状中的文本
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                
                if slide_text:
                    slides.append({
                        "text": f"--- 幻灯片 {i + 1} ---\n" + "\n".join(slide_text) + "\n",
                        "slide_number": i + 1
                    })
            
            return slides
        except Exception as e:
            print(f"加载PPTX文件出错: {file_path}, 错误: {str(e)}")
            return []

    def load_docx(self, file_path: str) -> str:
        """加载DOCX文件"""
        try:
            # 使用docx2txt提取文本
            text = docx2txt.process(file_path)
            return text
        except Exception as e:
            print(f"加载DOCX文件出错: {file_path}, 错误: {str(e)}")
            return ""

    def load_txt(self, file_path: str) -> str:
        """加载TXT文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # 尝试其他编码
            try:
                with open(file_path, 'r', encoding='gbk') as f:
                    return f.read()
            except Exception as e:
                print(f"加载TXT文件出错（编码问题）: {file_path}, 错误: {str(e)}")
                return ""
        except Exception as e:
            print(f"加载TXT文件出错: {file_path}, 错误: {str(e)}")
            return ""

    def load_document(self, file_path: str) -> List[Dict[str, str]]:
        """加载单个文档，PDF和PPT按页/幻灯片分割，返回文档块列表"""
        ext = os.path.splitext(file_path)[1].lower()
        filename = os.path.basename(file_path)
        documents = []

        if ext == ".pdf":
            pages = self.load_pdf(file_path)
            for page_idx, page_data in enumerate(pages, 1):
                documents.append(
                    {
                        "content": page_data["text"],
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": page_idx,
                    }
                )
        elif ext == ".pptx":
            slides = self.load_pptx(file_path)
            for slide_idx, slide_data in enumerate(slides, 1):
                documents.append(
                    {
                        "content": slide_data["text"],
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": slide_idx,
                    }
                )
        elif ext == ".docx":
            content = self.load_docx(file_path)
            if content:
                documents.append(
                    {
                        "content": content,
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": 0,
                    }
                )
        elif ext == ".txt":
            content = self.load_txt(file_path)
            if content:
                documents.append(
                    {
                        "content": content,
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": 0,
                    }
                )
        else:
            print(f"不支持的文件格式: {ext}")

        return documents

    def load_all_documents(self) -> List[Dict[str, str]]:
        """加载数据目录下的所有文档"""
        if not os.path.exists(self.data_dir):
            print(f"数据目录不存在: {self.data_dir}")
            return []

        documents = []

        for root, dirs, files in os.walk(self.data_dir):
            for file in files:
                ext = os.path.splitext(file)[1].lower()
                if ext in self.supported_formats:
                    file_path = os.path.join(root, file)
                    print(f"正在加载: {file_path}")
                    doc_chunks = self.load_document(file_path)
                    if doc_chunks:
                        documents.extend(doc_chunks)

        print(f"共加载 {len(documents)} 个文档块")
        return documents